import math
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor


MAJOR_LEVEL_VALUE = "3 - Major"


def format_date_for_ppt(value):
    if pd.isna(value) or value is None:
        return "N/A"
    try:
        dt = pd.to_datetime(value)
        return dt.strftime("%d/%m/%Y")
    except Exception:
        return str(value)


def safe_text(value, default="N/A"):
    if value is None or pd.isna(value):
        return default
    text = str(value).strip()
    return text if text else default


def clean_title(value):
    text = safe_text(value, default="N/A")
    if text.upper().startswith("BP2I-"):
        return text[5:].strip()
    return text


def normalize_tag(value):
    if value is None or pd.isna(value):
        return ""
    return str(value).strip()


def normalize_response(value):
    if value is None or pd.isna(value):
        return ""
    return str(value).strip().lower()


def normalize_major(value):
    if value is None or pd.isna(value):
        return ""
    return str(value).strip().lower()


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=12,
    bold=False,
    font_color=RGBColor(0, 0, 0),
    align=PP_ALIGN.LEFT,
    vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    rotation=0,
    margin_left=0,
    margin_right=0,
    margin_top=0,
    margin_bottom=0,
):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.rotation = rotation

    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = vertical_anchor
    text_frame.margin_left = margin_left
    text_frame.margin_right = margin_right
    text_frame.margin_top = margin_top
    text_frame.margin_bottom = margin_bottom

    p = text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = font_color

    return textbox


def filter_major_rows(df, column_name="Residual risk level"):
    if column_name not in df.columns:
        raise ValueError(f"Colonne manquante pour le filtre major: {column_name}")

    target = normalize_major(MAJOR_LEVEL_VALUE)
    return df[df[column_name].apply(normalize_major) == target].copy()


def build_rl_to_tag_map(risk_cards_df):
    """
    Construit un mapping RL -> TAG à partir de l'extract risk cards.
    On prend le premier TAG non vide trouvé pour chaque RL.
    """
    possible_rl_columns = [
        "Local risk reference",
        "Local Risk Reference",
        "RL",
        "Risk ID",
        "Risk Id",
    ]

    rl_col = None
    for col in possible_rl_columns:
        if col in risk_cards_df.columns:
            rl_col = col
            break

    if rl_col is None:
        return {}

    temp = risk_cards_df[[rl_col, "TAG"]].copy()
    temp[rl_col] = temp[rl_col].astype(str).str.strip()
    temp["TAG"] = temp["TAG"].apply(normalize_tag)

    temp = temp[(temp[rl_col] != "") & (temp["TAG"] != "")]
    if temp.empty:
        return {}

    temp = temp.drop_duplicates(subset=[rl_col], keep="first")

    return dict(zip(temp[rl_col], temp["TAG"]))


def compute_metrics(
    ipt_input_file="all_ipt.xlsx",
    risk_cards_input_file="risk_cards_extract.xlsx",
    excel_output="ipt_metrics.xlsx",
    ppt_output="ipt_progress_report.pptx",
):
    # =========================
    # Lecture IPT
    # =========================
    ipt_all = pd.read_excel(ipt_input_file)

    expected_ipt_columns = [
        "State",
        "Local risk reference",
        "Percent Complete",
        "Planned end date",
        "Risk Owner/Sponsor",
        "Residual risk level",
        "Title",
    ]
    missing_ipt = [c for c in expected_ipt_columns if c not in ipt_all.columns]
    if missing_ipt:
        raise ValueError(
            f"Colonnes manquantes dans le fichier IPT source: {missing_ipt}"
        )

    ipt_all["Percent Complete"] = pd.to_numeric(
        ipt_all["Percent Complete"], errors="coerce"
    ).fillna(0)

    ipt_all["Planned end date"] = pd.to_datetime(
        ipt_all["Planned end date"], errors="coerce"
    )

    # Filtre Major dès le départ
    ipt_all = filter_major_rows(ipt_all, column_name="Residual risk level")

    # =========================
    # Lecture Risk Cards
    # =========================
    risk_cards = pd.read_excel(risk_cards_input_file)

    expected_rc_columns = ["Response", "Title", "TAG", "Residual risk level"]
    missing_rc = [c for c in expected_rc_columns if c not in risk_cards.columns]
    if missing_rc:
        raise ValueError(
            f"Colonnes manquantes dans le fichier Risk Cards: {missing_rc}"
        )

    # Filtre Major dès le départ
    risk_cards = filter_major_rows(risk_cards, column_name="Residual risk level")

    # Mapping RL -> TAG pour classer aussi les IPT
    rl_to_tag = build_rl_to_tag_map(risk_cards)

    # =========================
    # Split IPT open / closed
    # =========================
    ipt_closed = ipt_all[ipt_all["State"] == "Closed Complete"].copy()
    ipt_open = ipt_all[ipt_all["State"] != "Closed Complete"].copy()

    rls_closed = set(
        ipt_closed["Local risk reference"].dropna().astype(str).str.strip().unique()
    )
    rls_open = set(
        ipt_open["Local risk reference"].dropna().astype(str).str.strip().unique()
    )
    all_rls = sorted(rls_closed.union(rls_open))

    results = []

    # =========================
    # Construction métriques IPT
    # =========================
    for rl in all_rls:
        closed_subset = ipt_closed[
            ipt_closed["Local risk reference"].astype(str).str.strip() == rl
        ]
        open_subset = ipt_open[
            ipt_open["Local risk reference"].astype(str).str.strip() == rl
        ]

        closed_count = closed_subset.shape[0]
        open_count = open_subset.shape[0]
        total = closed_count + open_count

        completions_open = open_subset["Percent Complete"].tolist()
        completions_closed = [100] * closed_count
        all_completions = completions_open + completions_closed
        avg_completion = (
            round(sum(all_completions) / len(all_completions), 2)
            if all_completions
            else 0
        )

        ratio_open = f"{open_count}/{total}" if total > 0 else "0/0"

        planned_dates = open_subset["Planned end date"].dropna()
        latest_planned_end_date = (
            planned_dates.max() if not planned_dates.empty else None
        )

        sample_row = pd.concat([open_subset, closed_subset], ignore_index=True).head(1)

        risk_owner = safe_text(
            sample_row["Risk Owner/Sponsor"].iloc[0] if not sample_row.empty else None
        )
        residual_level = safe_text(
            sample_row["Residual risk level"].iloc[0] if not sample_row.empty else None
        )
        title = clean_title(
            sample_row["Title"].iloc[0] if not sample_row.empty else rl
        )

        tag = rl_to_tag.get(str(rl).strip(), "")

        results.append(
            {
                "Source Type": "IPT",
                "Local risk reference": rl,
                "Title": title,
                "Open actions ratio": ratio_open,
                "Average completion rate (%)": avg_completion,
                "Latest planned end date": latest_planned_end_date,
                "Risk Owner/Sponsor": risk_owner,
                "Residual risk level": residual_level,
                "Bar Label": f"{avg_completion:.2f}".rstrip("0").rstrip(".") + "%",
                "Is Acceptance": False,
                "TAG": tag,
            }
        )

    # =========================
    # Ajout des Risk Cards Accept (Major uniquement déjà filtré)
    # =========================
    risk_cards_accept = risk_cards[
        risk_cards["Response"].apply(normalize_response) == "accept"
    ].copy()

    for _, row in risk_cards_accept.iterrows():
        title = clean_title(row["Title"])
        tag = normalize_tag(row["TAG"])

        results.append(
            {
                "Source Type": "Risk Card Acceptance",
                "Local risk reference": "N/A",
                "Title": title,
                "Open actions ratio": "N/A",
                "Average completion rate (%)": 0,
                "Latest planned end date": None,
                "Risk Owner/Sponsor": "N/A",
                "Residual risk level": safe_text(row.get("Residual risk level")),
                "Bar Label": "Acceptance",
                "Is Acceptance": True,
                "TAG": tag,
            }
        )

    df_results = pd.DataFrame(results)

    with pd.ExcelWriter(excel_output, engine="openpyxl") as writer:
        df_results.to_excel(writer, sheet_name="Metrics", index=False)

    create_progress_ppt(df_results, ppt_output)

    print(f"Excel généré : {excel_output}")
    print(f"PowerPoint généré : {ppt_output}")

    return df_results


def draw_group_frame(slide):
    outer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.25),
        Inches(0.95),
        Inches(12.75),
        Inches(5.95),
    )
    outer.fill.background()
    outer.line.color.rgb = RGBColor(0, 99, 92)
    outer.line.width = Pt(3)

    top_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.25),
        Inches(0.95),
        Inches(12.75),
        Inches(0.08),
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = RGBColor(0, 99, 92)
    top_line.line.fill.background()


def add_vertical_tag_panel(slide, tag_text):
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.35),
        Inches(1.08),
        Inches(0.72),
        Inches(5.68),
    )
    panel.fill.background()
    panel.line.color.rgb = RGBColor(0, 99, 92)
    panel.line.width = Pt(2)

    add_textbox(
        slide,
        Inches(0.40),
        Inches(2.05),
        Inches(0.55),
        Inches(3.5),
        tag_text,
        font_size=16,
        bold=True,
        font_color=RGBColor(0, 99, 92),
        align=PP_ALIGN.CENTER,
        rotation=270,
    )


def add_slide_header(slide, section_name=None):
    add_textbox(
        slide,
        Inches(10.9),
        Inches(0.72),
        Inches(1.5),
        Inches(0.22),
        "actions",
        font_size=11,
        bold=False,
        font_color=RGBColor(80, 80, 80),
        align=PP_ALIGN.CENTER,
    )

    draw_group_frame(slide)

    if section_name:
        add_vertical_tag_panel(slide, section_name)


def add_row_to_slide(slide, row, row_index_in_slide, has_tag_panel):
    text_dark = RGBColor(0, 0, 0)
    text_light = RGBColor(255, 255, 255)
    green = RGBColor(67, 160, 71)
    bar_bg = RGBColor(255, 255, 255)
    border = RGBColor(60, 60, 60)
    amber = RGBColor(193, 134, 0)

    current_top = Inches(1.18) + row_index_in_slide * Inches(0.43)

    if has_tag_panel:
        label_left = Inches(1.18)
        label_width = Inches(8.35)
    else:
        label_left = Inches(0.45)
        label_width = Inches(9.05)

    ratio_left = label_left + label_width + Inches(0.12)
    ratio_width = Inches(0.62)

    bar_left = ratio_left + ratio_width + Inches(0.12)
    bar_width = Inches(1.32)
    bar_height = Inches(0.29)

    date_left = bar_left + bar_width + Inches(0.14)
    date_width = Inches(1.02)

    title_text = clean_title(row.get("Title", "N/A"))
    progress = float(row.get("Average completion rate (%)", 0) or 0)
    progress = max(0, min(100, progress))
    ratio_text = safe_text(row.get("Open actions ratio", "N/A"))
    planned_end_date = format_date_for_ppt(row.get("Latest planned end date"))
    bar_label = safe_text(row.get("Bar Label", ""))
    is_acceptance = bool(row.get("Is Acceptance", False))

    title_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        label_left,
        current_top,
        label_width,
        Inches(0.33),
    )
    title_box.fill.background()
    title_box.line.color.rgb = border
    title_box.line.width = Pt(1.3)

    add_textbox(
        slide,
        label_left + Inches(0.05),
        current_top + Inches(0.01),
        label_width - Inches(0.08),
        Inches(0.27),
        title_text,
        font_size=10,
        bold=False,
        font_color=text_dark,
        align=PP_ALIGN.LEFT,
    )

    add_textbox(
        slide,
        ratio_left,
        current_top,
        ratio_width,
        Inches(0.33),
        ratio_text,
        font_size=10,
        bold=False,
        font_color=text_dark,
        align=PP_ALIGN.CENTER,
    )

    bg_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        bar_left,
        current_top + Inches(0.02),
        bar_width,
        bar_height,
    )
    bg_bar.fill.solid()
    bg_bar.fill.fore_color.rgb = bar_bg
    bg_bar.line.color.rgb = border
    bg_bar.line.width = Pt(1.1)
    bg_bar.line.dash_style = 1

    if not is_acceptance and progress > 0:
        fill_width = bar_width * (progress / 100.0)
        fg_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            bar_left,
            current_top + Inches(0.02),
            fill_width,
            bar_height,
        )
        fg_bar.fill.solid()
        fg_bar.fill.fore_color.rgb = green
        fg_bar.line.fill.background()

    add_textbox(
        slide,
        bar_left,
        current_top + Inches(0.01),
        bar_width,
        Inches(0.24),
        bar_label,
        font_size=9,
        bold=False,
        font_color=text_dark if is_acceptance else (text_light if progress >= 35 else text_dark),
        align=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        date_left,
        current_top,
        date_width,
        Inches(0.33),
        planned_end_date,
        font_size=10,
        bold=False,
        font_color=amber if planned_end_date != "N/A" else text_dark,
        align=PP_ALIGN.LEFT,
    )


def create_progress_ppt(df, output_file):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    rows_per_slide = 13

    df = df.copy()
    df["TAG"] = df["TAG"].fillna("").astype(str).str.strip()

    tagged_df = df[df["TAG"] != ""].copy()
    untagged_df = df[df["TAG"] == ""].copy()

    source_order_tagged = {"IPT": 0, "Risk Card Acceptance": 1}
    tagged_df["__source_order"] = tagged_df["Source Type"].map(source_order_tagged).fillna(99)
    tagged_df = (
        tagged_df.sort_values(
            by=["TAG", "__source_order", "Title"],
            ascending=[True, True, True]
        )
        .drop(columns="__source_order")
        .reset_index(drop=True)
    )

    source_order_untagged = {"IPT": 0, "Risk Card Acceptance": 1}
    untagged_df["__source_order"] = untagged_df["Source Type"].map(source_order_untagged).fillna(99)
    untagged_df = (
        untagged_df.sort_values(
            by=["__source_order", "Title"],
            ascending=[True, True]
        )
        .drop(columns="__source_order")
        .reset_index(drop=True)
    )

    if not tagged_df.empty:
        for tag_value, group_df in tagged_df.groupby("TAG", sort=True):
            group_df = group_df.reset_index(drop=True)
            total_rows = len(group_df)
            total_slides = math.ceil(total_rows / rows_per_slide)

            for slide_part in range(total_slides):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_slide_header(slide, section_name=tag_value)

                start = slide_part * rows_per_slide
                end = min(start + rows_per_slide, total_rows)
                chunk = group_df.iloc[start:end]

                for i, (_, row) in enumerate(chunk.iterrows()):
                    add_row_to_slide(
                        slide=slide,
                        row=row,
                        row_index_in_slide=i,
                        has_tag_panel=True,
                    )

    if not untagged_df.empty:
        total_rows = len(untagged_df)
        total_slides = math.ceil(total_rows / rows_per_slide)

        for slide_part in range(total_slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_slide_header(slide, section_name=None)

            start = slide_part * rows_per_slide
            end = min(start + rows_per_slide, total_rows)
            chunk = untagged_df.iloc[start:end]

            for i, (_, row) in enumerate(chunk.iterrows()):
                add_row_to_slide(
                    slide=slide,
                    row=row,
                    row_index_in_slide=i,
                    has_tag_panel=False,
                )

    if len(prs.slides) == 0:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_textbox(
            slide,
            Inches(1),
            Inches(3),
            Inches(6),
            Inches(0.5),
            "No major risk data available",
            font_size=18,
            bold=True,
        )

    prs.save(output_file)


if __name__ == "__main__":
    compute_metrics(
        ipt_input_file="all_ipt.xlsx",
        risk_cards_input_file="risk_cards_extract.xlsx",
        excel_output="ipt_metrics.xlsx",
        ppt_output="ipt_progress_report.pptx",
    )